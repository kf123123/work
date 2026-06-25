# -*- coding: utf-8 -*-
"""
PPT 生成脚本 — 车辆行驶远程控制系统设计
时尚科技风格，与中期报告对应
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============ 路径 ============
RESOURCE_DIR = "D:/UserData/Desktop/practice/resource"
OUTPUT_DIR = "D:/UserData/Desktop/practice/report"
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "车辆行驶远程控制系统_中期答辩PPT.pptx")

# ============ 配色方案 ============
COLOR_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 深蓝黑背景
COLOR_BG_CARD = RGBColor(0x1F, 0x29, 0x37)        # 卡片灰
COLOR_TEAL = RGBColor(0x2D, 0xD4, 0xBF)           # 青绿色（主色调）
COLOR_TEAL_LIGHT = RGBColor(0x5E, 0xEA, 0xD4)     # 浅青绿
COLOR_WHITE = RGBColor(0xE2, 0xE8, 0xF0)          # 白色文字
COLOR_GRAY = RGBColor(0x9C, 0xA3, 0xAF)           # 灰色
COLOR_ACCENT = RGBColor(0x60, 0xA5, 0xFA)         # 蓝色强调
COLOR_ORANGE = RGBColor(0xFB, 0x92, 0x3C)         # 橙色强调
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)            # 红色

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============ 工具函数 ============

def add_bg(slide, color=COLOR_BG_DARK):
    """设置幻灯片背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        # set alpha on the solidFill
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16,
                       color=COLOR_WHITE, font_name="Microsoft YaHei"):
    """带项目符号的文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_side_accent(slide, color=COLOR_TEAL):
    """左侧装饰条纹"""
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, color)

def add_top_bar(slide, text, color=COLOR_TEAL):
    """顶部标题栏"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), color)
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.7),
                text, font_size=28, color=color, bold=True)
    # 底部细线
    add_rect(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.02), COLOR_GRAY)

def add_image_centered(slide, img_path, left=None, top=None, width=None, height=None):
    """居中插入图片"""
    if not os.path.exists(img_path):
        add_textbox(slide, Inches(2), Inches(2), Inches(8), Inches(1),
                    f"[图片未找到: {os.path.basename(img_path)}]",
                    font_size=14, color=COLOR_RED)
        return
    if left is None:
        left = Inches(1.5)
    if top is None:
        top = Inches(1.5)
    if width is None:
        width = Inches(10)
    try:
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception:
        add_textbox(slide, Inches(2), Inches(2), Inches(8), Inches(1),
                    f"[图片加载失败]", font_size=14, color=COLOR_RED)


# ============================================================
#  Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 左上装饰
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, COLOR_TEAL)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), COLOR_TEAL)

# 主标题
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "车辆行驶远程控制系统", font_size=44, color=COLOR_TEAL, bold=True)

# 副标题
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
            "中期答辩报告", font_size=28, color=COLOR_WHITE, bold=False)

# 分隔线
add_rect(slide, Inches(1.5), Inches(3.7), Inches(3), Inches(0.04), COLOR_TEAL)

# 信息
info_items = [
    "专    业：电子科学与技术",
    "指导老师：王红亮",
    "组    别：第三组",
    "答辩日期：2026年6月",
]
for i, item in enumerate(info_items):
    add_textbox(slide, Inches(1.5), Inches(4.1 + i * 0.5), Inches(6), Inches(0.5),
                item, font_size=18, color=COLOR_GRAY)

# 组员信息右侧卡片
add_rect(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(4.5), COLOR_BG_CARD)
add_textbox(slide, Inches(8.3), Inches(1.7), Inches(4), Inches(0.5),
            "项目组成员", font_size=20, color=COLOR_TEAL, bold=True)

members = [
    ("2306024123", "李兆新", "UWB定位解析 / 文档"),
    ("2306024131", "李科锋", "上位机 / 视频传输 / 协议"),
    ("2306024132", "白  宇", "电机驱动 / 下位机协议"),
]
y = 2.4
for sid, name, role in members:
    add_textbox(slide, Inches(8.3), Inches(y), Inches(4), Inches(0.4),
                f"{name}  {sid}", font_size=18, color=COLOR_WHITE, bold=True)
    add_textbox(slide, Inches(8.3), Inches(y + 0.35), Inches(4), Inches(0.3),
                f"  {role}", font_size=13, color=COLOR_GRAY)
    y += 0.95


# ============================================================
#  Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "目  录  CONTENTS")

toc = [
    ("01", "项目概述", "背景与任务要求"),
    ("02", "系统总体方案", "架构设计 / 通信方案 / 硬件选型"),
    ("03", "通信协议定义", "控制指令 / 定位数据 / 视频传输协议"),
    ("04", "硬件设计与实现", "电路接口 / 电机驱动原理"),
    ("05", "软件设计与实现", "上位机 / 下位机 / 关键代码"),
    ("06", "定位模块", "UWB原理 / 协议解析 / 现状说明"),
    ("07", "实验与测试", "测试方法 / 测试结果 / 现场照片"),
    ("08", "总结与展望", "完成情况 / 后续计划"),
]

x_start = Inches(1)
y_start = Inches(1.3)
for i, (num, title, desc) in enumerate(toc):
    col = i % 4
    row = i // 4
    x = x_start + col * Inches(3)
    y = y_start + row * Inches(2.8)

    # 卡片背景
    add_rect(slide, x, y, Inches(2.7), Inches(2.3), COLOR_BG_CARD)

    # 编号
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(0.6),
                num, font_size=32, color=COLOR_TEAL, bold=True)

    # 标题
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.3), Inches(0.5),
                title, font_size=18, color=COLOR_WHITE, bold=True)

    # 描述
    add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.3), Inches(0.7),
                desc, font_size=11, color=COLOR_GRAY)


# ============================================================
#  Slide 3: 项目概述
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "一、项目概述")

# 左栏
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "项目背景", font_size=22, color=COLOR_TEAL, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5),
            "随着物联网技术与无线通信的快速发展，远程控制系统在智能交通、物流配送等领域得到广泛应用。本项目设计并实现一套基于Wi-Fi的车辆远程控制系统，实现上位机对车辆的实时视频监控与运动控制。",
            font_size=15, color=COLOR_WHITE)

add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
            "任务书要求", font_size=22, color=COLOR_TEAL, bold=True)
items = [
    "▸  上位机监控界面设计（PyQt5）",
    "▸  实时视频传输（UDP分片）",
    "▸  车辆运动控制（8方向 + 调速）",
    "▸  UWB定位追踪",
]
add_bullet_textbox(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5),
                   items, font_size=15, color=COLOR_WHITE)

# 右栏 - 技术栈
add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(7.6), Inches(1.5), Inches(4.5), Inches(0.5),
            "技术栈", font_size=22, color=COLOR_TEAL, bold=True)

tech_items = [
    ("上位机", "PyQt5 + OpenCV + NumPy"),
    ("视频传输", "K230 MicroPython UDP"),
    ("主控MCU", "ESP32 Arduino"),
    ("电机驱动", "PWM H桥差速控制"),
    ("定位模块", "Nooploop UWB NLink"),
    ("通信协议", "TCP JSON + UDP"),
]
y = 2.3
for label, val in tech_items:
    add_textbox(slide, Inches(7.6), Inches(y), Inches(2), Inches(0.35),
                label, font_size=14, color=COLOR_ACCENT, bold=True)
    add_textbox(slide, Inches(9.5), Inches(y), Inches(3), Inches(0.35),
                val, font_size=14, color=COLOR_WHITE)
    y += 0.65


# ============================================================
#  Slide 4: 系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "二、系统总体方案")

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.5),
            "系统架构设计", font_size=22, color=COLOR_TEAL, bold=True)

# 架构图 - 三个模块
modules = [
    ("上位机 (PC)", "PyQt5监控中心", [
        "视频画面显示",
        "8方向控制面板",
        "定位地图与轨迹",
        "状态信息显示",
    ], Inches(0.8), Inches(2.0)),
    ("Wi-Fi 网络", "UDP 8080 + TCP 8888", [
        "视频流 UDP 分片",
        "控制指令 TCP JSON",
        "定位数据 TCP JSON",
    ], Inches(4.8), Inches(2.0)),
    ("下位机 (K230 + ESP32)", "硬件控制层", [
        "K230 摄像头推流",
        "ESP32 电机驱动",
        "UWB 定位数据采集",
    ], Inches(8.8), Inches(2.0)),
]

for title, subtitle, items, x, y in modules:
    # 模块卡片
    add_rect(slide, x, y, Inches(3.6), Inches(4.5), COLOR_BG_CARD)
    # 标题
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.2), Inches(0.5),
                title, font_size=20, color=COLOR_TEAL, bold=True)
    # 副标题
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.2), Inches(0.4),
                subtitle, font_size=12, color=COLOR_ACCENT)
    # 分隔
    add_rect(slide, x + Inches(0.2), y + Inches(1.15), Inches(3.2), Inches(0.02), COLOR_TEAL)
    # 列表
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(1.4 + j * 0.55), Inches(3.0), Inches(0.5),
                    f"◆  {item}", font_size=14, color=COLOR_WHITE)

    # 箭头
    if title != "下位机 (K230 + ESP32)":
        arrow_x = x + Inches(3.6)
        add_textbox(slide, arrow_x + Inches(0.1), y + Inches(1.8), Inches(0.8), Inches(0.5),
                    "⟷", font_size=28, color=COLOR_TEAL, bold=True)

# 底部说明
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
            "通信方式：视频 → UDP（端口8080） | 控制指令 → TCP（端口8888） | 定位数据 → TCP（端口8888）",
            font_size=13, color=COLOR_GRAY)


# ============================================================
#  Slide 5: 通信协议
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "三、通信协议定义")

# 三种协议卡片
protocols = [
    ("控制指令协议", "PC → ESP32", "TCP JSON", [
        "T=11: PWM控制(L/R)",
        "T=115: 停止",
        "T=116~123: 方向控制",
        "T=200/201: UWB控制",
        "T=999: 心跳",
    ]),
    ("定位数据协议", "ESP32 → PC", "TCP JSON", [
        "T=200 固定消息类型",
        "x/y: 车辆坐标(米)",
        "valid: 数据有效性",
        "t: 时间戳(ms)",
        "推送频率: 5Hz",
    ]),
    ("视频传输协议", "K230 → PC", "UDP分片", [
        "每帧1400字节分片",
        "6字节包头: frame_id",
        "total_chunks / chunk_id",
        "JPEG payload",
        "quality=35优化",
    ]),
]

x = Inches(0.6)
for title, direction, protocol, items in protocols:
    add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(5.2), COLOR_BG_CARD)

    # 协议名标题
    add_textbox(slide, x + Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.5),
                title, font_size=22, color=COLOR_TEAL, bold=True)

    # 方向标签
    add_rect(slide, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.35), RGBColor(0x37, 0x41, 0x51))
    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(2), Inches(0.35),
                direction, font_size=13, color=COLOR_ACCENT)
    # 协议类型
    add_textbox(slide, x + Inches(2.2), Inches(2.1), Inches(1.5), Inches(0.35),
                protocol, font_size=13, color=COLOR_TEAL_LIGHT, bold=True)

    # 分隔
    add_rect(slide, x + Inches(0.2), Inches(2.6), Inches(3.4), Inches(0.02), COLOR_TEAL)

    # 内容
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), Inches(2.9 + j * 0.5), Inches(3.3), Inches(0.45),
                    item, font_size=14, color=COLOR_WHITE)

    # 代码示例区域
    if title == "控制指令协议":
        code_text = '{"T":116,"pwm":150}'
    elif title == "定位数据协议":
        code_text = '{"T":200,"x":1.23,"y":0.56}'
    else:
        code_text = "[ID:2][总:2][片:1][JPEG]"

    add_rect(slide, x + Inches(0.2), Inches(5.5), Inches(3.4), Inches(0.6), RGBColor(0x11, 0x18, 0x27))
    add_textbox(slide, x + Inches(0.4), Inches(5.55), Inches(3.2), Inches(0.5),
                code_text, font_size=12, color=COLOR_TEAL_LIGHT)

    x += Inches(4.1)


# ============================================================
#  Slide 6: 硬件设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "四、硬件设计与实现")

# 左 - 引脚表
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
            "ESP32 电路接口", font_size=22, color=COLOR_TEAL, bold=True)

pins = [
    "PIN_PWMA (GPIO25) — 左电机PWM",
    "PIN_AIN1 (GPIO21) — 左电机方向",
    "PIN_AIN2 (GPIO17) — 左电机方向",
    "PIN_BIN1 (GPIO22) — 右电机方向",
    "PIN_BIN2 (GPIO23) — 右电机方向",
    "PIN_PWMB (GPIO26) — 右电机PWM",
    "UART2 RX (GPIO16) — UWB数据接收",
]
for i, pin in enumerate(pins):
    add_textbox(slide, Inches(0.8), Inches(2.0 + i * 0.5), Inches(5.5), Inches(0.45),
                f"  {pin}", font_size=13, color=COLOR_WHITE if i < 6 else COLOR_ORANGE)

# 右 - 运动控制表
add_textbox(slide, Inches(7), Inches(1.3), Inches(5), Inches(0.5),
            "差速驱动方式", font_size=22, color=COLOR_TEAL, bold=True)

motions = [
    ("前进", "正转", "正转"),
    ("后退", "反转", "反转"),
    ("左转", "反转", "正转"),
    ("右转", "正转", "反转"),
    ("左前转", "停止", "正转"),
    ("右前转", "正转", "停止"),
    ("停止", "停止", "停止"),
]

# 表头
add_rect(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(0.45), RGBColor(0x2D, 0xD4, 0xBF))
add_textbox(slide, Inches(7.1), Inches(1.92), Inches(1.8), Inches(0.4),
            "运动状态", font_size=14, color=RGBColor(0x11, 0x18, 0x27), bold=True)
add_textbox(slide, Inches(9.0), Inches(1.92), Inches(1.5), Inches(0.4),
            "左轮", font_size=14, color=RGBColor(0x11, 0x18, 0x27), bold=True)
add_textbox(slide, Inches(10.8), Inches(1.92), Inches(1.5), Inches(0.4),
            "右轮", font_size=14, color=RGBColor(0x11, 0x18, 0x27), bold=True)

for i, (state, left, right) in enumerate(motions):
    y = Inches(2.4 + i * 0.45)
    bg = COLOR_BG_CARD if i % 2 == 0 else RGBColor(0x1A, 0x23, 0x32)
    add_rect(slide, Inches(7), y, Inches(5.5), Inches(0.42), bg)
    add_textbox(slide, Inches(7.1), y + Inches(0.02), Inches(1.8), Inches(0.38),
                state, font_size=13, color=COLOR_WHITE, bold=(state == "停止"))
    add_textbox(slide, Inches(9.0), y + Inches(0.02), Inches(1.5), Inches(0.38),
                left, font_size=13, color=COLOR_WHITE)
    add_textbox(slide, Inches(10.8), y + Inches(0.02), Inches(1.5), Inches(0.38),
                right, font_size=13, color=COLOR_WHITE)


# ============================================================
#  Slide 7: 软件设计 - 上位机
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "五、软件设计与实现 — 上位机")

# 三个功能模块卡片
modules_ui = [
    ("VideoReceiveThread", "视频接收线程", [
        "监听UDP 8080端口",
        "JPEG分片接收与组包",
        "OpenCV解码显示",
        "帧率统计",
    ]),
    ("TcpReceiveThread", "定位接收线程", [
        "TCP JSON逐行解析",
        "UWB坐标提取",
        "信号发射到主界面",
        "断线重连处理",
    ]),
    ("ControlWindow", "主控窗口", [
        "8方向按钮 + 键盘控制",
        "速度滑块实时调参",
        "定位地图与轨迹",
        "状态信息表格",
    ]),
]

x = Inches(0.6)
for title, subtitle, items in modules_ui:
    add_rect(slide, x, Inches(1.3), Inches(3.8), Inches(3.2), COLOR_BG_CARD)

    add_textbox(slide, x + Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.5),
                title, font_size=20, color=COLOR_TEAL_LIGHT, bold=True)
    add_textbox(slide, x + Inches(0.2), Inches(1.95), Inches(3.4), Inches(0.3),
                subtitle, font_size=12, color=COLOR_GRAY)
    add_rect(slide, x + Inches(0.2), Inches(2.3), Inches(3.4), Inches(0.02), COLOR_TEAL)

    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), Inches(2.5 + j * 0.5), Inches(3.2), Inches(0.45),
                    f"◆  {item}", font_size=14, color=COLOR_WHITE)
    x += Inches(4.1)

# 界面截图
if os.path.exists(os.path.join(RESOURCE_DIR, "pyqt.jpg")):
    add_image_centered(slide,
        os.path.join(RESOURCE_DIR, "pyqt.jpg"),
        left=Inches(1.5), top=Inches(4.8), width=Inches(10), height=None)
    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(5), Inches(0.4),
                "▲ 上位机监控主界面", font_size=13, color=COLOR_GRAY)

# 控制方式说明
add_textbox(slide, Inches(0.8), Inches(6.9), Inches(12), Inches(0.4),
            "控制方式: 鼠标点击D-pad按钮  |  键盘 WASD移动 + 空格停止 + Q/E/Z/C对角  |  速度滑块 0~200",
            font_size=13, color=COLOR_ACCENT)


# ============================================================
#  Slide 8: 软件设计 - 下位机
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "五、软件设计与实现 — 下位机")

# 左 - ESP32
add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.5),
            "ESP32 固件 (Arduino)", font_size=22, color=COLOR_TEAL, bold=True)
add_rect(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(0.02), COLOR_TEAL)

esp_items = [
    "主程序: setup() 初始化 → loop() 循环执行",
    "├─ chassisInit() — 电机引脚初始化",
    "├─ initUwbSerial() — UWB串口初始化",
    "├─ wifiTcpInit() — WiFi + TCP服务启动",
    "├─ uwbSerialRead() — 高频读取UWB数据",
    "├─ wifiTcpLoop() — TCP处理+指令解析+UWB推送",
    "└─ handleSerial() — 串口调试命令处理",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(5.4), Inches(3.5),
                   esp_items, font_size=14, color=COLOR_WHITE)

# 电机驱动逻辑
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5.4), Inches(0.4),
            "电机驱动: H桥控制左右轮，PWM调速，差速转向",
            font_size=14, color=COLOR_ACCENT)

# 右 - K230
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.5),
            "K230 摄像头推流 (MicroPython)", font_size=22, color=COLOR_TEAL, bold=True)
add_rect(slide, Inches(7.0), Inches(2.1), Inches(5.4), Inches(0.02), COLOR_TEAL)

k230_items = [
    "Wi-Fi联网 → 连接同一局域网",
    "摄像头初始化 (640×480 RGB565)",
    "JPEG压缩 (quality=35优化帧率)",
    "UDP分片发送 (每片≤1400字节)",
    "帧ID + 分片序号组包恢复",
]
add_bullet_textbox(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(3.0),
                   k230_items, font_size=14, color=COLOR_WHITE)

# 关键代码
add_textbox(slide, Inches(7.0), Inches(4.5), Inches(5.4), Inches(0.4),
            "关键代码示意:", font_size=15, color=COLOR_TEAL, bold=True)
code_bg = add_rect(slide, Inches(7.0), Inches(5.0), Inches(5.4), Inches(1.5), RGBColor(0x11, 0x18, 0x27))
code_text = "struct.pack('!HHH', frame_id, total_chunks, chunk_id)\nudp_socket.sendto(header + payload, dest_addr)"
add_textbox(slide, Inches(7.2), Inches(5.1), Inches(5.0), Inches(1.3),
            code_text, font_size=12, color=COLOR_TEAL_LIGHT)


# ============================================================
#  Slide 9: UWB定位
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "六、定位模块设计与实现")

# 左
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "UWB 定位原理", font_size=22, color=COLOR_TEAL, bold=True)
uwb_items = [
    "◆ 基于TDOA到达时间差定位",
    "◆ Nooploop LinkTrack 模块",
    "◆ 精度: 10~30cm (视环境)",
    "◆ 通过UART2串口接收数据",
    "◆ 波特率 921600 高速模式",
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(2.5),
                   uwb_items, font_size=15, color=COLOR_WHITE)

add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
            "NLink 协议解析", font_size=22, color=COLOR_TEAL, bold=True)
nlink_items = [
    "▸ 帧同步: 0x55起始字节",
    "▸ Tag_Frame0 (0x01): 10字节定长帧",
    "▸ Node_Frame2 (0x04): 可变长帧",
    "▸ int24编码坐标 → 毫米→米转换",
    "▸ 校验和验证数据完整性",
]
add_bullet_textbox(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5),
                   nlink_items, font_size=14, color=COLOR_WHITE)

# 右 - 现状
add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(7.5), Inches(1.5), Inches(4.8), Inches(0.5),
            "模块现状", font_size=22, color=COLOR_ORANGE, bold=True)
add_rect(slide, Inches(7.5), Inches(2.1), Inches(4.8), Inches(0.02), COLOR_ORANGE)

status_items = [
    "✅ 协议解析代码已完成",
    "✅ 模拟器验证通过",
    "✅ 上位机地图显示正常",
    "",
    "⚠️ 硬件模块已损坏",
    "⚠️ 后续更换模块后联调",
]
add_bullet_textbox(slide, Inches(7.5), Inches(2.3), Inches(4.8), Inches(3.0),
                   status_items, font_size=16, color=COLOR_WHITE)

# 数据格式
add_textbox(slide, Inches(7.5), Inches(5.0), Inches(4.8), Inches(0.4),
            "定位数据 JSON 格式:", font_size=15, color=COLOR_TEAL, bold=True)
code_bg = add_rect(slide, Inches(7.5), Inches(5.5), Inches(4.8), Inches(0.9), RGBColor(0x11, 0x18, 0x27))
add_textbox(slide, Inches(7.7), Inches(5.6), Inches(4.4), Inches(0.7),
            '{"T":200,"x":1.23,"y":0.56,"valid":true,"t":12345}',
            font_size=13, color=COLOR_TEAL_LIGHT)


# ============================================================
#  Slide 10: 实验与测试
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "七、实验与测试")

# 测试项目卡片
tests = [
    ("上位机功能测试", "✅ 通过", [
        "8方向按钮正常响应",
        "键盘控制WASD正常",
        "速度滑块实时调参",
    ]),
    ("视频传输测试", "✅ 通过", [
        "本地模拟 25~30 FPS",
        "硬件实测 ~20 FPS",
        "UDP分片组包正确",
    ]),
    ("电机控制测试", "✅ 通过", [
        "9种运动状态正常",
        "PWM调速功能正常",
        "正反转控制正确",
    ]),
    ("UWB定位测试", "⚠️ 部分通过", [
        "模拟器验证通过",
        "地图显示与轨迹正常",
        "硬件模块待更换",
    ]),
]

x = Inches(0.6)
for title, status, items in tests:
    add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(3.5), COLOR_BG_CARD)
    add_textbox(slide, x + Inches(0.2), Inches(1.5), Inches(2.5), Inches(0.4),
                title, font_size=16, color=COLOR_TEAL, bold=True)
    status_color = COLOR_TEAL if "通过" in status else COLOR_ORANGE
    add_textbox(slide, x + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.35),
                status, font_size=13, color=status_color, bold=True)
    add_rect(slide, x + Inches(0.2), Inches(2.3), Inches(2.5), Inches(0.02), COLOR_GRAY)
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), Inches(2.5 + j * 0.45), Inches(2.5), Inches(0.4),
                    f"• {item}", font_size=12, color=COLOR_WHITE)
    x += Inches(3.15)

# 测试照片
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(5), Inches(0.4),
            "测试现场照片", font_size=20, color=COLOR_TEAL, bold=True)

if os.path.exists(os.path.join(RESOURCE_DIR, "practice1.jpg")):
    add_image_centered(slide,
        os.path.join(RESOURCE_DIR, "practice1.jpg"),
        left=Inches(0.8), top=Inches(5.6), width=Inches(5.5), height=None)

if os.path.exists(os.path.join(RESOURCE_DIR, "practice2.jpg")):
    add_image_centered(slide,
        os.path.join(RESOURCE_DIR, "practice2.jpg"),
        left=Inches(7.0), top=Inches(5.6), width=Inches(5.5), height=None)


# ============================================================
#  Slide 11: 完成情况
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "八、完成情况分析")

# 完成度仪表盘
add_rect(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2), COLOR_BG_CARD)
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.0),
            "整体完成度评估：核心功能已基本实现，UWB定位模块代码已完成但硬件损坏待更换",
            font_size=18, color=COLOR_ORANGE, bold=True)

# 完成情况表
headers = ["要求", "状态", "说明"]
rows_data = [
    ("上位机界面", "✅ 已完成", "PyQt5开发，视频/控制/状态/地图四区域"),
    ("视频传输", "✅ 已完成", "UDP分片传输，JPEG quality=35"),
    ("运动控制", "✅ 已完成", "8方向+速度调节，按钮/键盘双模式"),
    ("UWB定位", "⚠️ 部分", "协议解析代码完成，硬件损坏待换"),
    ("系统联调", "🔄 基本完成", "各模块单元测试通过"),
]

# 表头
add_rect(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5), COLOR_TEAL)
for ci, h in enumerate(headers):
    xs = [Inches(0.8), Inches(4.5), Inches(7.5)]
    add_textbox(slide, xs[ci], Inches(2.83), Inches(3), Inches(0.4),
                h, font_size=15, color=RGBColor(0x11, 0x18, 0x27), bold=True)

for ri, (req, status, desc) in enumerate(rows_data):
    y = Inches(3.4 + ri * 0.5)
    bg = COLOR_BG_CARD if ri % 2 == 0 else RGBColor(0x1A, 0x23, 0x32)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(0.45), bg)
    add_textbox(slide, Inches(0.8), y + Inches(0.03), Inches(3.5), Inches(0.38),
                req, font_size=14, color=COLOR_WHITE, bold=True)
    add_textbox(slide, Inches(4.5), y + Inches(0.03), Inches(2.8), Inches(0.38),
                status, font_size=14, color=COLOR_TEAL if "完成" in status else COLOR_ORANGE)
    add_textbox(slide, Inches(7.5), y + Inches(0.03), Inches(4.8), Inches(0.38),
                desc, font_size=13, color=COLOR_GRAY)

# 底部总结
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
            "未完成项分析：UWB定位模块已完成完整的代码开发和模拟器验证，但实际硬件模块在联调阶段故障。\n"
            "后续计划采购同型号模块替换，完成定位功能的最终验证。",
            font_size=14, color=COLOR_ACCENT)


# ============================================================
#  Slide 12: 总结与展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_side_accent(slide)
add_top_bar(slide, "九、总结与展望")

# 左 - 已完成
add_rect(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.5),
            "已完成工作", font_size=22, color=COLOR_TEAL, bold=True)
add_rect(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(0.02), COLOR_TEAL)

done_items = [
    "✅ PyQt5 上位机监控软件开发",
    "✅ K230 摄像头UDP推流程序",
    "✅ ESP32 电机驱动 + WiFi + UWB固件",
    "✅ 完整的通信协议栈",
    "✅ 系统联调测试",
    "✅ 项目文档与报告",
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(5.4), Inches(3.5),
                   done_items, font_size=15, color=COLOR_WHITE)

# 右 - 后续
add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), COLOR_BG_CARD)
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.4), Inches(0.5),
            "后续工作计划", font_size=22, color=COLOR_ORANGE, bold=True)
add_rect(slide, Inches(7.0), Inches(2.1), Inches(5.4), Inches(0.02), COLOR_ORANGE)

plan_items = [
    "➊ 更换UWB模块，完成定位联调",
    "➋ 优化视频帧率与传输稳定性",
    "➌ 引入PID控制优化电机精度",
    "➍ 完善系统文档与测试数据",
    "➎ 结题验收准备",
]
add_bullet_textbox(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(3.0),
                   plan_items, font_size=16, color=COLOR_WHITE)

# 感谢
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.8),
            "感谢各位老师聆听！",
            font_size=28, color=COLOR_TEAL, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
#  保存
# ============================================================
prs.save(OUTPUT_FILE)
print(f"PPT已生成: {OUTPUT_FILE}")
